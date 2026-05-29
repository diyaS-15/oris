# Parses uploaded PPTX and PDF files into a list of chunks, one per slide or page.
# Strips repeating headers (any first line appearing in more than 3 chunks) before storing.
import io
from collections import Counter
from dataclasses import dataclass

import pdfplumber
from pptx import Presentation


@dataclass
class Chunk:
    content: str
    source_file: str
    source_type: str          # 'slide' | 'practice'
    page_or_slide_number: int


def parse_file(filename: str, file_bytes: bytes, source_type: str) -> list[Chunk]:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pptx":
        return _parse_pptx(filename, file_bytes, source_type)
    elif ext == "pdf":
        return _parse_pdf(filename, file_bytes, source_type)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _strip_repeating_headers(chunks: list[Chunk]) -> list[Chunk]:
    """
    Detects lines that appear at the start of more than 3 chunks in the same file
    and strips them from every chunk. Handles stacked headers (e.g. course name on
    line 1, professor name on line 2) by iterating until no more are found.
    """
    result = [c for c in chunks]  # shallow copy so we don't mutate originals

    while True:
        # Count how many chunks start with each distinct first line
        first_line_counts: Counter = Counter()
        for chunk in result:
            first_line = chunk.content.split("\n")[0].strip()
            if first_line:
                first_line_counts[first_line] += 1

        # Any line that appears at the top of more than 3 chunks is a repeating header
        headers = {line for line, count in first_line_counts.items() if count > 3}
        if not headers:
            break

        stripped = []
        for chunk in result:
            lines = chunk.content.split("\n")
            # Remove all leading lines that are repeating headers
            while lines and lines[0].strip() in headers:
                lines = lines[1:]
            new_content = "\n".join(lines).strip()
            if new_content:
                stripped.append(Chunk(
                    content=new_content,
                    source_file=chunk.source_file,
                    source_type=chunk.source_type,
                    page_or_slide_number=chunk.page_or_slide_number,
                ))
            # If stripping leaves the chunk empty, drop it entirely

        result = stripped

    return result


def _parse_pptx(filename: str, file_bytes: bytes, source_type: str) -> list[Chunk]:
    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        content = "\n".join(texts)
        if content:
            chunks.append(Chunk(content=content, source_file=filename,
                                source_type=source_type, page_or_slide_number=i))
    return _strip_repeating_headers(chunks)


def _parse_pdf(filename: str, file_bytes: bytes, source_type: str) -> list[Chunk]:
    chunks = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            content = (page.extract_text() or "").strip()
            if content:
                chunks.append(Chunk(content=content, source_file=filename,
                                    source_type=source_type, page_or_slide_number=i))
    return _strip_repeating_headers(chunks)
